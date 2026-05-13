import logging
import requests
from io import BytesIO
from typing import Tuple, Optional

logger = logging.getLogger(__name__)


def extract_text(
    source_type: str,
    source_url: Optional[str] = None,
    raw_text: Optional[str] = None,
) -> Tuple[str, Optional[int]]:
    """
    Extract text from various sources.
    Returns (text, page_count)
    """
    try:
        if source_type == "raw_text":
            return (raw_text or ""), None

        if source_type == "url":
            return _extract_url(source_url), None

        if source_type == "image":
            return _extract_image(source_url), None

        # File-based: download from Cloudinary URL
        if not source_url:
            raise ValueError("source_url required for file types")

        response = requests.get(source_url, timeout=30)

        if response.status_code == 401:
            raise ValueError(
                "File could not be accessed (401 Unauthorized). "
                "In Cloudinary Dashboard → Settings → Security, ensure "
                "'Strict delivery mode' is disabled, or make the file publicly accessible. "
                f"URL attempted: {source_url}"
            )
        elif response.status_code == 403:
            raise ValueError(
                "File access forbidden (403). The Cloudinary file may have restricted "
                "delivery settings. Check your Cloudinary account's access control settings."
            )
        elif response.status_code == 404:
            raise ValueError(
                f"File not found (404). The file may have been deleted from Cloudinary. "
                f"URL: {source_url}"
            )
        elif response.status_code == 500:
            raise ValueError(
                "Cloudinary server error (500). This may indicate account restrictions "
                f"or misconfiguration. URL: {source_url}"
            )

        response.raise_for_status()
        buffer = BytesIO(response.content)

        if source_type == "pdf":
            return _extract_pdf(buffer)
        elif source_type == "docx":
            return _extract_docx(buffer), None
        elif source_type == "ppt":
            return _extract_pptx(buffer), None
        elif source_type == "txt":
            return buffer.read().decode("utf-8", errors="ignore").strip(), None
        else:
            raise ValueError(f"Unsupported source type: {source_type}")

    except Exception as e:
        logger.error(f"[extractor] Error extracting {source_type}: {e}")
        raise


def _extract_pdf(buffer: BytesIO) -> Tuple[str, int]:
    from pypdf import PdfReader
    reader = PdfReader(buffer)
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        pages.append(text)
    return "\n\n".join(pages), len(reader.pages)


def _extract_docx(buffer: BytesIO) -> str:
    from docx import Document
    doc = Document(buffer)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_pptx(buffer: BytesIO) -> str:
    from pptx import Presentation
    prs = Presentation(buffer)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides_text.append(f"[Slide {i+1}]\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


def _extract_url(url: str) -> str:
    from bs4 import BeautifulSoup
    import re

    # ✅ FIX: Pehle check karo reddit hai ya nahi, phir properly replace karo
    # Bug tha: "reddit.com" replace hone se "old.old.reddit.com" ban raha tha
    is_reddit = "reddit.com" in url
    if is_reddit:
        # www.reddit.com → old.reddit.com, ya reddit.com → old.reddit.com
        # re.sub use karo taake sirf ek baar replace ho
        fetch_url = re.sub(r'(https?://)(?:www\.)?reddit\.com', r'\1old.reddit.com', url)
    else:
        fetch_url = url

    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/120.0.0.0 Safari/537.36"
        ),
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.5",
    }

    response = requests.get(fetch_url, headers=headers, timeout=30)
    response.raise_for_status()

    soup = BeautifulSoup(response.text, "html.parser")

    for tag in soup(["script", "style", "nav", "footer", "header",
                     "aside", "form", "noscript", "iframe"]):
        tag.decompose()

    if is_reddit:
        content_el = (
            soup.find("div", class_="thing")
            or soup.find("div", class_="entry")
            or soup.find("div", id="siteTable")
            or soup.find("div", class_="commentarea")
        )
    else:
        content_el = (
            soup.find("article")
            or soup.find("main")
            or soup.find(id="content")
            or soup.find(class_=lambda c: c and any(
                x in c for x in ["content", "article", "post", "entry"]
            ))
            or soup.find("body")
        )

    text = content_el.get_text(separator="\n") if content_el else soup.get_text(separator="\n")

    lines = [line.strip() for line in text.splitlines()]
    lines = [line for line in lines if len(line) > 20]
    result = "\n\n".join(lines[:500])

    if not result.strip():
        body = soup.find("body")
        if body:
            all_lines = [l.strip() for l in body.get_text(separator="\n").splitlines()]
            all_lines = [l for l in all_lines if len(l) > 15]
            result = "\n\n".join(all_lines[:500])

    if not result.strip():
        raise ValueError(
            f"No readable text could be extracted from URL: {url}. "
            "The page may be JavaScript-rendered, paywalled, or require login. "
            "Try copying the text manually and using the 'Text' input instead."
        )

    return result

def _extract_image(url: str) -> str:
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        raise ValueError(
            "Image text extraction requires pytesseract and Pillow. "
            "Install with: pip install pytesseract pillow\n"
            "Also install Tesseract OCR: https://github.com/tesseract-ocr/tesseract"
        )

    response = requests.get(url, timeout=30)
    response.raise_for_status()
    img = Image.open(BytesIO(response.content))
    text = pytesseract.image_to_string(img).strip()

    if not text:
        raise ValueError(
            "No text could be extracted from the image. "
            "The image may not contain readable text, or the text may be too small/blurry."
        )
    return text